"""PPTX 文本抽取（python-pptx）。

遍历每张幻灯片的 shapes（文本框、表格、占位符），按幻灯片顺序输出文本。
"""

from __future__ import annotations

import asyncio
import logging
from pathlib import Path

from ..models import ExtractedContent

logger = logging.getLogger(__name__)


async def extract(path: Path, max_chars: int = 0) -> ExtractedContent:
    """从 PPTX 抽取纯文本。"""
    return await asyncio.to_thread(_extract_sync, path, max_chars)


def _extract_sync(path: Path, max_chars: int) -> ExtractedContent:
    try:
        from pptx import Presentation
    except ImportError:
        logger.error("python-pptx not installed; cannot extract .pptx")
        return ExtractedContent(text="", file_type="pptx", metadata={"error": "python-pptx not installed"})

    prs = Presentation(str(path))
    slides_count = len(prs.slides)
    parts: list[str] = []
    total = 0

    for idx, slide in enumerate(prs.slides, start=1):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            # 文本框 / 占位符
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        slide_texts.append(line)
            # 表格
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    line = " | ".join(cells)
                    if line.strip(" |"):
                        slide_texts.append(line)

        if slide_texts:
            parts.append(f"--- Slide {idx} ---")
            parts.extend(slide_texts)
            total += sum(len(s) for s in slide_texts)
            if max_chars > 0 and total >= max_chars:
                text = "\n".join(parts)[:max_chars]
                return ExtractedContent(
                    text=text, file_type="pptx", pages=slides_count,
                    char_count=len(text), metadata={"truncated": True},
                )

    text = "\n".join(parts)
    return ExtractedContent(
        text=text,
        file_type="pptx",
        pages=slides_count,
        char_count=len(text),
    )
